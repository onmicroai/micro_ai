import os
import pytesseract
import pandas as pd
import fitz  # PyMuPDF for PDFs
import pdfplumber
import docx
import csv
from PIL import Image
from pptx import Presentation
from abc import ABC, abstractmethod

# Abstract Class for Document Parsers
class DocumentParser(ABC):
    """Abstract base class for document parsers."""
    @abstractmethod
    def extract_text(self, file_path):
        pass

# PDF Parser (Handles Normal and Scanned PDFs)
class PDFParser(DocumentParser):
    def extract_text(self, file_path):
        text = ""
        try:
            doc = fitz.open(file_path)
            text = "\n".join([page.get_text("text") for page in doc])
        except Exception:
            pass
        if not text.strip():
            try:
                with pdfplumber.open(file_path) as pdf:
                    text = "\n".join([page.extract_text() or "" for page in pdf.pages])
            except Exception:
                pass
        if not text.strip():
            text = OCRParser().extract_text(file_path)  # Use OCR for scanned PDFs
        return text.strip()

# Word Document Parser
class WordParser(DocumentParser):
    def extract_text(self, file_path):
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs]).strip()

# Excel Parser (Handles XLSX, CSV)
class ExcelParser(DocumentParser):
    def extract_text(self, file_path):
        text = []
        if file_path.endswith(".csv"):
            with open(file_path, "r", encoding="utf-8") as file:
                reader = csv.reader(file)
                text = ["\t".join(row) for row in reader]
        else:
            df = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, sheet_data in df.items():
                text.append(f"--- Sheet: {sheet_name} ---")
                text.append(sheet_data.to_string(index=False, header=True))
        return "\n".join(text).strip()

# Text File Parser
class TextParser(DocumentParser):
    def extract_text(self, file_path):
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read().strip()

# PowerPoint Parser
class PowerPointParser(DocumentParser):
    def extract_text(self, file_path):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text).strip()

# OCR Parser (Handles Images and Scanned PDFs)
class OCRParser(DocumentParser):
    def extract_text(self, file_path):
        image = Image.open(file_path)
        return pytesseract.image_to_string(image).strip()

# Context Class for Document Processing and Validation
class DocumentProcessor:
    """Determines the appropriate parser for a document and validates file type and size."""

    MAX_FILE_SIZE_MB = 10  # 10MB limit

    def __init__(self):
        self.parsers = {
            ".pdf": PDFParser(),
            ".docx": WordParser(),
            ".doc": WordParser(),
            ".xlsx": ExcelParser(),
            ".xls": ExcelParser(),
            ".csv": ExcelParser(),
            ".txt": TextParser(),
            ".log": TextParser(),
            ".pptx": PowerPointParser(),
            ".jpg": OCRParser(),
            ".jpeg": OCRParser(),
            ".png": OCRParser(),
            ".tiff": OCRParser(),
            ".bmp": OCRParser(),
        }

    def validate_file(self, file_path):
        """Checks if file type is supported and size is within limit."""
        ext = os.path.splitext(file_path)[1].lower()
        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)  # Convert to MB

        if ext not in self.parsers:
            return f"Error: Unsupported file format ({ext})"

        if file_size_mb > self.MAX_FILE_SIZE_MB:
            return f"Error: File size exceeds 10MB limit ({file_size_mb:.2f} MB)"

        return "valid"

    def extract_text(self, file_path):
        """Validates and extracts text from the file."""
        validation_result = self.validate_file(file_path)
        if validation_result != "valid":
            return validation_result

        ext = os.path.splitext(file_path)[1].lower()
        parser = self.parsers.get(ext)

        if parser:
            return parser.extract_text(file_path)
        else:
            return "Unsupported file format"